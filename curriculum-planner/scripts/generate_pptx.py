#!/usr/bin/env python3
"""
Generate branded PPTX presentations for each session by copying the ALC
template and modifying slides in-place.

Usage:
    python scripts/generate_pptx.py

Reads:  output/parsed_curriculum.json
        templates/ALC_Curriculum_Template.pptx
Output: output/Session_1_Slides.pptx (etc.)
"""

import json
import shutil
import sys
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Pt

SCRIPT_DIR = Path(__file__).parent.parent
OUTPUT_DIR = SCRIPT_DIR / "output"
TEMPLATE_PATH = SCRIPT_DIR / "templates" / "ALC_Curriculum_Template.pptx"
PARSED_PATH = OUTPUT_DIR / "parsed_curriculum.json"

# Template slide indices (0-based)
SLIDE_TITLE = 0
SLIDE_MODULE = 1
SLIDE_OBJECTIVES = 2
SLIDE_VOCABULARY = 3
SLIDE_ACTIVITY = 4
SLIDE_ASSESSMENT = 5
SLIDE_CLOSING = 6


def set_shape_text(shape, text):
    """Replace ALL text in a shape while preserving first-run formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        for run in tf.paragraphs[0].runs[1:]:
            run.text = ""
        for para in tf.paragraphs[1:]:
            for run in para.runs:
                run.text = ""
    else:
        tf.text = text


def duplicate_slide(prs, source_index):
    """Duplicate an existing slide in the presentation. Returns the new slide."""
    source = prs.slides[source_index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Remove default shapes
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy all shapes from source
    for shape in source.shapes:
        new_slide.shapes._spTree.append(deepcopy(shape._element))

    return new_slide


def build_session_pptx(session, output_path):
    """Build a PPTX for one session by copying template and modifying in-place."""
    # Copy template file
    shutil.copy2(str(TEMPLATE_PATH), str(output_path))
    prs = Presentation(str(output_path))

    num = session["session_num"]
    title = session["title"]
    subtitle = session["subtitle"]
    schedule = session["schedule"]
    swbats = session["swbats"]
    activities = [b for b in schedule if not b.get("is_break")]

    # Parse subtitle parts
    parts = subtitle.split("|") if subtitle else ["", "", ""]
    day_part = parts[0].strip() if len(parts) > 0 else ""
    time_part = parts[1].strip() if len(parts) > 1 else ""

    # === SLIDE 1: TITLE ===
    slide = prs.slides[SLIDE_TITLE]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        tl = text.lower()
        if "business english" in tl:
            paras = shape.text_frame.paragraphs
            if len(paras) >= 2 and paras[0].runs and paras[1].runs:
                paras[0].runs[0].text = "AI & Digital Workflows"
                paras[1].runs[0].text = "Bootcamp"
        elif "b1 intermediate" in tl or "contact hours" in tl:
            set_shape_text(shape, f"Session {num}: {title}  |  7 Contact Hours")

    # === SLIDE 2: MODULE HEADER ===
    slide = prs.slides[SLIDE_MODULE]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "MODULE" in text:
            set_shape_text(shape, f"SESSION {num:02d}")
        elif "FUTURE SKILLS" in text or "INNOVATION" in text:
            set_shape_text(shape, title.upper())
        elif "Sessions" in text or "Hours" in text:
            set_shape_text(shape, f"{day_part}  |  {time_part}")

    # === SLIDE 3: LEARNING OBJECTIVES ===
    slide = prs.slides[SLIDE_OBJECTIVES]
    clean_swbats = [s.split("(LO")[0].strip().rstrip(" ;,") for s in swbats]

    obj_texts = []
    bullet_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text == "▸":
            bullet_shapes.append(shape)
        elif "Module" in text and len(text) < 60:
            set_shape_text(shape, f"Session {num}: {title}")
        elif len(text) > 10 and 900000 < shape.top < 4800000:
            obj_texts.append(shape)

    obj_texts.sort(key=lambda s: s.top)
    bullet_shapes.sort(key=lambda s: s.top)
    for i, shape in enumerate(obj_texts):
        set_shape_text(shape, clean_swbats[i] if i < len(clean_swbats) else "")
    for i, shape in enumerate(bullet_shapes):
        if i >= len(clean_swbats):
            set_shape_text(shape, "")

    # === SLIDE 4: KEY VOCABULARY ===
    slide = prs.slides[SLIDE_VOCABULARY]
    vocab = get_session_vocabulary(num)

    left_terms = []
    right_contexts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text in ("KEY VOCABULARY", "TERMS", "IN CONTEXT", "ALC"):
            continue
        if "Module" in text and len(text) < 60:
            set_shape_text(shape, f"Session {num}: {title}")
            continue
        if shape.top > 1000000 and shape.top < 4800000:
            if shape.left < 4500000:
                left_terms.append(shape)
            else:
                right_contexts.append(shape)

    left_terms.sort(key=lambda s: s.top)
    right_contexts.sort(key=lambda s: s.top)

    for i, shape in enumerate(left_terms):
        if i < len(vocab):
            paras = shape.text_frame.paragraphs
            if len(paras) >= 2 and paras[0].runs and paras[1].runs:
                paras[0].runs[0].text = vocab[i][0]
                paras[1].runs[0].text = vocab[i][1]
            else:
                set_shape_text(shape, f"{vocab[i][0]}\n{vocab[i][1]}")
        else:
            set_shape_text(shape, "")

    for i, shape in enumerate(right_contexts):
        if i < len(vocab) and len(vocab[i]) > 2:
            set_shape_text(shape, f'"{vocab[i][2]}"')
        else:
            set_shape_text(shape, "")

    # === SLIDE 5: FIRST ACTIVITY (modify in-place) ===
    major_activities = [b for b in activities if int(b["duration"].replace(" min", "")) >= 25]
    if major_activities:
        populate_activity_slide(prs.slides[SLIDE_ACTIVITY], major_activities[0], session)

    # === SLIDE 6: ASSESSMENT ===
    slide = prs.slides[SLIDE_ASSESSMENT]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "Module" in text and len(text) < 60:
                set_shape_text(shape, f"Session {num}: {title}")
        if shape.has_table:
            criteria = build_assessment_criteria(num)
            for ri in range(1, len(shape.table.rows)):
                if ri - 1 < len(criteria):
                    for ci, cell in enumerate(shape.table.rows[ri].cells):
                        if ci < len(criteria[ri - 1]):
                            for para in cell.text_frame.paragraphs:
                                if para.runs:
                                    para.runs[0].text = criteria[ri - 1][ci]
                                else:
                                    para.text = criteria[ri - 1][ci]

    # === SLIDE 7: CLOSING ===
    next_titles = {
        1: "Session 2: Research, Write, Create",
        2: "Session 3: Automate & Strategize",
        3: "Session 4: Ship It & Own It",
        4: "Bootcamp Complete — Build Your Future with AI",
    }
    slide = prs.slides[SLIDE_CLOSING]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "Module 02" in text or "Writing Slack" in text:
                set_shape_text(shape, next_titles.get(num, ""))

    # === ADD EXTRA ACTIVITY SLIDES ===
    # Duplicate the activity slide (now modified) for additional activities
    extra_activities = major_activities[1:6]  # up to 5 more
    for block in extra_activities:
        new_slide = duplicate_slide(prs, SLIDE_ACTIVITY)
        populate_activity_slide(new_slide, block, session)

    prs.save(str(output_path))
    return len(prs.slides)


def populate_activity_slide(slide, block, session):
    """Fill an activity slide with content from a schedule block."""
    num = session["session_num"]
    title = session["title"]
    steps = extract_steps(block.get("delivery", ""))

    # Collect all text shapes with their current text
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_shapes.append((shape, shape.text_frame.text.strip()))

    # Sort by position
    text_shapes.sort(key=lambda x: (x[0].top, x[0].left))

    # Identify shapes in the title zone (1M-1.8M top) with width > 5M
    title_zone = [(s, t) for s, t in text_shapes
                  if 1000000 <= s.top < 1800000 and s.width > 5000000 and t]
    title_zone.sort(key=lambda x: x[0].top)

    for shape, text in text_shapes:
        top = shape.top
        left = shape.left

        # Banner format/duration (top right)
        if top < 900000 and left > 3000000:
            fmt = block.get("format", "").split(".")[0].split("→")[0].strip()
            set_shape_text(shape, f"{fmt}  •  {block['duration']}")

        # Activity title area - first in zone = title, second = description
        elif 1000000 <= top < 1800000 and shape.width > 5000000:
            if title_zone and shape is title_zone[0][0]:
                set_shape_text(shape, block["activity"])
            elif len(title_zone) > 1 and shape is title_zone[1][0]:
                desc = block.get("delivery", "")[:120]
                if len(block.get("delivery", "")) > 120:
                    desc = desc.rsplit(" ", 1)[0] + "..."
                set_shape_text(shape, desc)

        # Steps area
        elif 2300000 <= top < 4300000:
            if text.isdigit():
                pass
            elif left > 1000000 and len(text) > 5:
                step_idx = None
                for other, other_text in text_shapes:
                    if (other_text.isdigit() and
                        abs(other.top - shape.top) < 50000):
                        step_idx = int(other_text) - 1
                        break
                if step_idx is not None and step_idx < len(steps):
                    set_shape_text(shape, steps[step_idx])
                elif step_idx is not None:
                    set_shape_text(shape, "")

        # Teacher tip
        elif 4300000 <= top < 4800000:
            tip = block.get("instructor_note", "")
            if not tip:
                tip = block.get("pacing", "")[:150] if block.get("pacing") else ""
            if tip:
                set_shape_text(shape, f"💡 Teacher tip: {tip[:200]}")
            else:
                set_shape_text(shape, "")

        # Footer
        elif top >= 4800000 and "ALC" not in text and "AMERICAN" not in text:
            set_shape_text(shape, f"Session {num}: {title}")


def extract_steps(delivery_text):
    """Extract actionable steps from a delivery description."""
    if not delivery_text:
        return []
    sentences = delivery_text.replace(". ", ".\n").split("\n")
    steps = []
    for s in sentences:
        s = s.strip()
        if s and len(s) > 15 and not s.startswith("This "):
            steps.append(s)
        if len(steps) >= 4:
            break
    return steps


def get_session_vocabulary(session_num):
    """Return vocabulary terms for each session."""
    return {
        1: [
            ("CRAFT", "Context, Role, Audience, Format, Tone",
             "Use CRAFT to structure every prompt for better output."),
            ("chain-of-thought", "Breaking tasks into step-by-step subtasks",
             "Think step by step: audience → messaging → channels."),
            ("token", "Smallest unit of text an LLM processes",
             "'Unbelievable' might be split into 3 tokens."),
            ("temperature", "Controls randomness (0=precise, 1=creative)",
             "Lower temperature for consistent, factual responses."),
        ],
        2: [
            ("hallucination", "AI generates plausible but fabricated info",
             "Always fact-check — AI confidently cites fake sources."),
            ("source triangulation", "Verifying across multiple tools/sources",
             "Same question through Claude, Perplexity, and NotebookLM."),
            ("AI collaboration", "Using AI at specific stages, not full replacement",
             "Best as brainstormer, outliner, or editor — not ghostwriter."),
            ("trust rubric", "Checklist for evaluating AI output reliability",
             "Cites sources? Verifiable? Holds up across tools?"),
        ],
        3: [
            ("trigger", "Event that starts an automation",
             "When I receive an email with 'invoice' in the subject…"),
            ("workflow automation", "Connecting tools for multi-step tasks",
             "Email trigger → AI extraction → spreadsheet update."),
            ("AI agent", "Program that plans, executes, iterates autonomously",
             "Reads your email, drafts replies, schedules follow-ups."),
            ("MCP", "Model Context Protocol — AI-to-software bridge",
             "Lets Claude talk to your calendar and documents."),
        ],
        4: [
            ("capstone", "Culminating project showing integrated skills",
             "Solve a real problem using multiple AI tools."),
            ("portfolio", "Curated collection of work artifacts",
             "Playbook, analysis, automation docs, and capstone."),
            ("Making Of", "Documentation of process and AI's role",
             "What did AI do well? What did you fix or override?"),
            ("peer review", "Structured evaluation using a rubric",
             "Does it solve a real problem? Is AI usage transparent?"),
        ],
    }.get(session_num, [])


def build_assessment_criteria(session_num):
    """Build assessment rubric rows for each session."""
    return {
        1: [
            ["Prompt quality", "Generic, no structure", "Uses CRAFT elements", "Polished, iterated prompts"],
            ["Tool awareness", "Can't distinguish tools", "Names differences", "Selects right tool for task"],
            ["Critical thinking", "Accepts all AI output", "Questions some output", "Identifies bias & errors"],
            ["Portfolio artifact", "Incomplete Playbook", "5 basic prompts", "7+ tested prompts w/ rubrics"],
        ],
        2: [
            ["Fact-checking", "Accepts AI claims", "Spots obvious errors", "Systematic verification"],
            ["Research skills", "Single-tool use", "Uses multiple tools", "Triangulates 3+ sources"],
            ["Writing voice", "Full AI replacement", "Some personal voice", "Clear human-AI collaboration"],
            ["Data analysis", "Raw AI output only", "Summarizes findings", "Insights w/ supporting evidence"],
        ],
        3: [
            ["Automation", "No working build", "Guided build works", "Independent build + iteration"],
            ["Workflow design", "No clear structure", "Basic trigger-action", "Multi-step w/ human checks"],
            ["Career positioning", "Vague aspirations", "Names AI skills", "Specific, actionable strategy"],
            ["Teaching ability", "Cannot explain", "Basic explanation", "Clear, accurate teach-back"],
        ],
        4: [
            ["Capstone quality", "Incomplete project", "Functional but basic", "Solves real problem"],
            ["Portfolio", "Missing artifacts", "All items present", "Curated w/ Making Of docs"],
            ["Peer review", "Surface-level", "Uses rubric criteria", "Actionable, specific feedback"],
            ["Action plan", "No commitments", "General goals", "Specific, time-bound actions"],
        ],
    }.get(session_num, [])


def main():
    OUTPUT_DIR.mkdir(exist_ok=True)

    if not PARSED_PATH.exists():
        print(f"Error: {PARSED_PATH} not found. Run parse_docx.py first.")
        sys.exit(1)

    if not TEMPLATE_PATH.exists():
        print(f"Error: Template not found: {TEMPLATE_PATH}")
        sys.exit(1)

    with open(PARSED_PATH) as f:
        curriculum = json.load(f)

    for session in curriculum["sessions"]:
        num = session["session_num"]
        filename = f"Session_{num}_Slides.pptx"
        output_path = OUTPUT_DIR / filename

        slide_count = build_session_pptx(session, output_path)
        print(f"Generated: {output_path} ({slide_count} slides)")

    print(f"\nAll {len(curriculum['sessions'])} presentations generated.")


if __name__ == "__main__":
    main()
