"""
Stage 4: GENERATE SLIDES — Create branded .pptx from lesson plans.

Reads parsed_curriculum.json and produces one .pptx per session in output/slides/.
Uses a branded template from templates/ if available, otherwise creates slides
from scratch.

Slide design rules:
- One idea per slide
- 6x6 ceiling: max 6 bullets, max 6 words per bullet
- No paragraphs on slides
- Slide types: Title, Objectives, Content, Activity, Discussion, CFU, Summary

Usage:
    python stage4_generate_slides.py
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


OUTPUT_DIR = Path(__file__).resolve().parent.parent / "output"
TEMPLATES_DIR = Path(__file__).resolve().parent.parent / "templates"
SLIDES_DIR = OUTPUT_DIR / "slides"

# Brand colors
COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)    # Blue
COLOR_SECONDARY = RGBColor(0x2D, 0x2D, 0x2D)  # Dark gray
COLOR_ACCENT = RGBColor(0xE8, 0x4D, 0x0E)      # Orange
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)


def load_parsed_curriculum() -> list[dict]:
    path = OUTPUT_DIR / "parsed_curriculum.json"
    if not path.exists():
        print(f"ERROR: {path} not found. Run stage1_parse.py first.")
        sys.exit(1)
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def get_template() -> Path | None:
    """Find a branded .pptx template in the templates directory."""
    templates = sorted(TEMPLATES_DIR.glob("*.pptx"))
    if templates:
        print(f"Using template: {templates[0].name}")
        return templates[0]
    return None


def _truncate_bullet(text: str, max_words: int = 6) -> str:
    """Enforce the 6-word-per-bullet ceiling."""
    words = text.split()
    if len(words) <= max_words:
        return text
    return " ".join(words[:max_words]) + "…"


def _add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Add a title slide."""
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def _add_content_slide(prs: Presentation, title: str, bullets: list[str],
                       max_bullets: int = 6):
    """Add a content slide with bullets (6x6 rule enforced)."""
    layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets[:max_bullets]):
        if i == 0:
            tf.text = _truncate_bullet(bullet)
        else:
            p = tf.add_paragraph()
            p.text = _truncate_bullet(bullet)
            p.level = 0


def _add_activity_slide(prs: Presentation, activity_name: str,
                        duration: int, format_type: str, instructions: str):
    """Add an activity slide."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = f"Activity: {activity_name}"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.text = f"Duration: {duration} min"
    p = tf.add_paragraph()
    p.text = f"Format: {format_type}"

    if instructions:
        steps = [s.strip() for s in instructions.split(".") if s.strip()]
        for step in steps[:5]:
            p = tf.add_paragraph()
            p.text = _truncate_bullet(step)
            p.level = 1


def _add_cfu_slide(prs: Presentation, question: str, technique: str = ""):
    """Add a check-for-understanding slide."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Check for Understanding"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    if technique:
        tf.text = f"Technique: {technique}"
        p = tf.add_paragraph()
        p.text = question
    else:
        tf.text = question


def generate_slides_for_session(session: dict):
    """Generate a .pptx file for one session."""
    num = session["session_number"]
    title = session["title"]
    swbats = session.get("swbats", [])
    blocks = session.get("schedule_blocks", [])

    template_path = get_template()
    if template_path:
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # Title slide
    _add_title_slide(prs, f"Session {num}: {title}")

    # Objectives slide
    if swbats:
        obj_bullets = [f"SWBAT: {s}" for s in swbats]
        _add_content_slide(prs, "Learning Objectives", obj_bullets)

    # Activity slides
    for block in blocks:
        name = block.get("activity_name", "Activity")
        duration = block.get("duration", 15)
        fmt = block.get("format", "Whole Class")
        delivery = block.get("delivery_instructions", "")

        _add_activity_slide(prs, name, duration, fmt, delivery)

    # CFU slide (placeholder)
    if swbats:
        _add_cfu_slide(prs, f"Can you demonstrate: {swbats[0]}?", "Think-Pair-Share")

    # Summary slide
    summary_bullets = [
        "Review key takeaways",
        "Complete exit ticket",
        "Preview next session",
    ]
    _add_content_slide(prs, "Summary & Closing", summary_bullets)

    # Save
    filename = f"session_{num:02d}_slides.pptx"
    out_path = SLIDES_DIR / filename
    prs.save(str(out_path))
    print(f"Generated: {filename}")


def main():
    sessions = load_parsed_curriculum()
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)

    for session in sessions:
        generate_slides_for_session(session)

    print(f"\n{len(sessions)} slide deck(s) written to {SLIDES_DIR}")


if __name__ == "__main__":
    main()
